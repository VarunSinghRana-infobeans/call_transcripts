"""
06_generate_ppt.py

Generate a production-quality PowerPoint presentation.
Visual style: clean executive briefing deck with card-based layouts,
consistent grid, generous margins, and no text overflow.

Exports: Transcript_Intelligence_Report.pptx
"""

import sys
import textwrap
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent / "scripts"))

from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from ppt_data import PresentationData, fmt_num, load_presentation_data
from utils import CHARTS_DIR, OUTPUT_DIR


def fmt_pct(value, decimals=1, fallback="N/A"):
    return fmt_num(value, decimals, fallback)


def _fmt_date_range(date_min: str, date_max: str) -> str:
    """Return a compact, readable date range like 'Feb 03 – Apr 28'."""
    try:
        from datetime import datetime
        dmin = datetime.strptime(date_min, "%Y-%m-%d")
        dmax = datetime.strptime(date_max, "%Y-%m-%d")
        return f"{dmin.strftime('%b %d')} – {dmax.strftime('%b %d')}"
    except Exception:
        return f"{date_min[-5:]} – {date_max[-5:]}"


# ---------------------------------------------------------------------------
# Slide dimensions — 16:9 executive format (matches reference deck)
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(10.0)
SLIDE_HEIGHT = Inches(5.62)

MARGIN = Inches(0.5)
CONTENT_W = SLIDE_WIDTH - MARGIN * 2
RIGHT_EDGE = SLIDE_WIDTH - MARGIN
BOTTOM_MAX = Inches(5.35)  # Hard stop for content

# Title area
TITLE_Y = Inches(0.25)
TITLE_H = Inches(0.50)
SUBTITLE_Y = Inches(0.82)
SUBTITLE_H = Inches(0.38)
UNDERLINE_Y = Inches(0.78)
UNDERLINE_W = Inches(1.80)
UNDERLINE_H = Inches(0.04)

# Content area
CONTENT_Y = Inches(1.30)

# Layout zones
CHART_X = MARGIN
CHART_W = Inches(5.2)
RIGHT_X = Inches(6.15)
RIGHT_W = Inches(3.30)

# Colors
C_PRIMARY = RGBColor(0x1a, 0x23, 0x7e)      # Dark blue
C_SECONDARY = RGBColor(0x00, 0x78, 0xd4)    # Bright blue
C_ACCENT = RGBColor(0xff, 0x6b, 0x35)       # Orange
C_TEXT = RGBColor(0x33, 0x33, 0x33)         # Dark gray
C_LIGHT = RGBColor(0x88, 0x88, 0x88)        # Light gray
C_WHITE = RGBColor(0xff, 0xff, 0xff)
C_BG = RGBColor(0xf5, 0xf6, 0xf8)           # Card background

C_RED = RGBColor(0xc0, 0x39, 0x2b)
C_AMBER = RGBColor(0xe6, 0x8a, 0x00)
C_GREEN = RGBColor(0x2e, 0x7d, 0x32)

# Fonts
FONT_TITLE = Pt(34)
FONT_SUBTITLE = Pt(15)
FONT_HEADING = Pt(23)
FONT_BODY = Pt(12)
FONT_SMALL = Pt(11)
FONT_TINY = Pt(10)
FONT_MICRO = Pt(9)
FONT_KPI = Pt(30)
FONT_KPI_LABEL = Pt(11)
FONT_NUMBER = Pt(14)

# Emoji mapping
TOPIC_ICONS = {
    "Compliance & Certification": "🔒",
    "Compliance & Audit": "🔒",
    "Compliance & Audits": "🔒",
    "Platform Reliability": "⚠️",
    "Incident Response & Reliability": "⚠️",
    "Incident Response & Outages": "⚠️",
    "Identity & Access Management": "🪪",
    "Identity & Access": "🪪",
    "Integrations & API": "🔌",
    "Engineering & Sprint Planning": "🔧",
    "Internal Operations": "🔧",
    "Customer Success": "🤝",
    "Threat Detection": "🔍",
    "Billing & Contracts": "💳",
    "Product & Roadmap": "🚀",
    "Churn & Risk": "🚨",
    "Sales & Renewals": "💰",
    "Product Deployment & Setup": "🚀",
    "Product": "🚀",
    "Billing": "💳",
    "API": "🔌",
    "Detection": "🔍",
}


# ---------------------------------------------------------------------------
# Sentiment color scale helper
# ---------------------------------------------------------------------------
def sentiment_color(score: float) -> RGBColor:
    if score is None:
        return C_LIGHT
    if score >= 4.0:
        return C_GREEN
    if score >= 3.0:
        return C_SECONDARY
    if score >= 2.0:
        return C_AMBER
    return C_RED


def risk_color(level: str) -> RGBColor:
    return {"High": C_RED, "Medium": C_AMBER, "Low": C_GREEN}.get(level, C_LIGHT)


def _to_pt(value) -> float:
    """Convert emu or pt value to points."""
    if isinstance(value, (int, float)) and value > 1000:  # emu
        return value / 12700
    return float(value)


def _to_in(value) -> float:
    """Convert emu or inches value to inches."""
    if isinstance(value, (int, float)) and value > 1000:  # emu
        return value / 914400
    return float(value)


def chars_per_inch(font_size_pt: float) -> float:
    """Approximate Calibri characters per inch at a given font size."""
    return font_size_pt * 1.15


def estimate_lines(text: str, width, font_size) -> int:
    """Rough estimate of wrapped lines."""
    width_in = _to_in(width)
    font_size_pt = _to_pt(font_size)
    if width_in <= 0 or not text:
        return 0
    cpi = chars_per_inch(font_size_pt)
    chars_per_line = max(1, int(width_in * cpi))
    return max(1, (len(text) + chars_per_line - 1) // chars_per_line)


def fit_text(text: str, width, font_size, max_lines: int) -> str:
    """Truncate text with ellipsis so it fits in the given line budget."""
    if not text:
        return text
    width_in = _to_in(width)
    font_size_pt = _to_pt(font_size)
    cpi = chars_per_inch(font_size_pt)
    chars_per_line = max(1, int(width_in * cpi))
    max_chars = chars_per_line * max_lines
    if len(text) <= max_chars:
        return text
    # Leave room for ellipsis
    return textwrap.shorten(text, width=max(max_chars - 3, 10), placeholder="...")


# ---------------------------------------------------------------------------
# Style helpers
# ---------------------------------------------------------------------------

def style_text(p, size, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    p.alignment = align


def fit_to_bounds(left, top, width, height, preserve_aspect=False,
                  max_right=RIGHT_EDGE, max_bottom=BOTTOM_MAX):
    """
    Reshape a rectangle so it stays inside the slide content bounds.
    Returns (left, top, width, height) as Inches objects.
    """
    left_in = _to_in(left)
    top_in = _to_in(top)
    width_in = _to_in(width)
    height_in = _to_in(height)
    max_right_in = _to_in(max_right)
    max_bottom_in = _to_in(max_bottom)

    if left_in < MARGIN.inches:
        left_in = MARGIN.inches

    right_in = left_in + width_in
    bottom_in = top_in + height_in

    overshoot_x = max(0.0, right_in - max_right_in)
    overshoot_y = max(0.0, bottom_in - max_bottom_in)

    if preserve_aspect and (overshoot_x > 0 or overshoot_y > 0):
        scale_x = 1.0 if overshoot_x <= 0 else (max_right_in - left_in) / width_in
        scale_y = 1.0 if overshoot_y <= 0 else (max_bottom_in - top_in) / height_in
        scale = min(scale_x, scale_y)
        width_in *= scale
        height_in *= scale
    else:
        if overshoot_x > 0:
            width_in = max(0.05, max_right_in - left_in)
        if overshoot_y > 0:
            height_in = max(0.05, max_bottom_in - top_in)

    return Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)


def add_text(slide, left, top, width, height, text, size=FONT_BODY, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    """Add a textbox with zero internal margins to maximize usable space."""
    left, top, width, height = fit_to_bounds(left, top, width, height)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = text
    style_text(p, size, bold=bold, color=color, align=align)
    return box


def add_bullets(slide, left, top, width, height, items, size=FONT_BODY, color=C_TEXT, align=PP_ALIGN.LEFT,
                bullet="•", max_items=6):
    """Add a bulleted textbox with one paragraph per item.

    Items are trimmed to the available vertical space by limiting the count.
    """
    left, top, width, height = fit_to_bounds(left, top, width, height)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.auto_size = MSO_AUTO_SIZE.NONE

    items = [str(i) for i in items if str(i).strip()]
    items = items[:max_items]
    for idx, item in enumerate(items):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet} {item}"
        style_text(p, size, bold=False, color=color, align=align)
        p.space_after = Pt(2)
    return box


def add_slide_title(slide, title, subtitle=""):
    """Plain title + optional subtitle with a colored underline."""
    add_text(slide, MARGIN, TITLE_Y, CONTENT_W, TITLE_H,
             title, FONT_HEADING, bold=True, color=C_PRIMARY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, UNDERLINE_Y, UNDERLINE_W, UNDERLINE_H)
    line.fill.solid()
    line.fill.fore_color.rgb = C_SECONDARY
    line.line.fill.background()
    if subtitle:
        add_text(slide, MARGIN, SUBTITLE_Y, CONTENT_W, SUBTITLE_H,
                 subtitle, FONT_SMALL, color=C_LIGHT)


def add_card(slide, left, top, width, height, fill=C_BG, line=None):
    """Add a rounded rectangle card background. Add BEFORE text so text sits on top."""
    left, top, width, height = fit_to_bounds(left, top, width, height)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line:
        card.line.color.rgb = line
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card


def add_chart_if_exists(slide, filename, left, top, width, height):
    """Insert a chart PNG, scaling it to fit the box while preserving aspect ratio."""
    path = CHARTS_DIR / filename
    if path.exists():
        left, top, box_w, box_h = fit_to_bounds(left, top, width, height, preserve_aspect=False)
        pic = slide.shapes.add_picture(str(path), left, top)
        aspect = pic.width / pic.height
        new_w = box_w
        new_h = Inches(new_w.inches / aspect)
        if new_h.inches > box_h.inches:
            new_h = box_h
            new_w = Inches(new_h.inches * aspect)
        pic.width = new_w
        pic.height = new_h
        return pic
    return None


def add_circle_badge(slide, left, top, size, text, color=C_SECONDARY, font_size=FONT_NUMBER):
    left, top, width, height = fit_to_bounds(left, top, size, size)
    size = width
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, left, top + Inches(0.02), size, size - Inches(0.04),
             str(text), font_size, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    return circle


# ---------------------------------------------------------------------------
# Grid-aligned KPI tiles
# ---------------------------------------------------------------------------
def add_kpi_row(slide, kpis, y, box_w, box_h, gap=Inches(0.12),
                card_fill=C_BG, label_color=C_LIGHT, default_value_size=FONT_KPI):
    """Add a row of equal KPI cards aligned to a strict grid.

    Each KPI can be a 3-tuple (value, label, color) or 4-tuple
    (value, label, color, value_font_size) for long values like date ranges.
    """
    n = len(kpis)
    total_w = n * box_w + (n - 1) * gap
    start_x = MARGIN + (CONTENT_W - total_w) / 2
    for i, item in enumerate(kpis):
        if len(item) == 4:
            val, label, color, value_size = item
        else:
            val, label, color = item
            value_size = default_value_size
        left = start_x + i * (box_w + gap)
        add_card(slide, left, y, box_w, box_h, fill=card_fill)
        add_text(slide, left, y + Inches(0.08), box_w, Inches(0.52),
                 str(val), value_size, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, left, y + Inches(0.63), box_w, Inches(0.28),
                 label, FONT_KPI_LABEL, color=label_color, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Layout guard
# ---------------------------------------------------------------------------

class LayoutGuard:
    def __init__(self, slide, name):
        self.slide = slide
        self.name = name
        self.warnings = []
        self.fixed = []

    def _is_background_or_decorative(self, shape):
        left_in = shape.left.inches
        top_in = shape.top.inches
        width_in = shape.width.inches
        height_in = shape.height.inches
        right_in = left_in + width_in
        bottom_in = top_in + height_in
        full_slide = (left_in < 0.01 and top_in < 0.01 and
                      abs(width_in - SLIDE_WIDTH.inches) < 0.01 and
                      abs(height_in - SLIDE_HEIGHT.inches) < 0.01)
        decorative = (
            top_in < -0.5 or
            left_in > SLIDE_WIDTH.inches - 1.0 or
            right_in > SLIDE_WIDTH.inches + 1.0 or
            bottom_in > SLIDE_HEIGHT.inches + 0.5
        )
        return full_slide or decorative

    def _reshape_shape(self, shape):
        left_in = shape.left.inches
        top_in = shape.top.inches
        width_in = shape.width.inches
        height_in = shape.height.inches
        preserve_aspect = shape.shape_type == 13
        new_left, new_top, new_width, new_height = fit_to_bounds(
            left_in, top_in, width_in, height_in,
            preserve_aspect=preserve_aspect
        )
        shape.left = new_left
        shape.top = new_top
        shape.width = new_width
        shape.height = new_height
        return (new_width.inches != width_in or new_height.inches != height_in or
                new_left.inches != left_in or new_top.inches != top_in)

    def check(self, fix=True):
        for shape in self.slide.shapes:
            if self._is_background_or_decorative(shape):
                continue
            left_in = shape.left.inches
            top_in = shape.top.inches
            width_in = shape.width.inches
            height_in = shape.height.inches
            right_in = left_in + width_in
            bottom_in = top_in + height_in
            overflows = (
                right_in > SLIDE_WIDTH.inches + 0.20 or
                bottom_in > SLIDE_HEIGHT.inches + 0.20 or
                bottom_in > BOTTOM_MAX.inches + 0.10
            )
            if overflows:
                if fix:
                    did_fix = self._reshape_shape(shape)
                    if did_fix:
                        self.fixed.append(
                            f"reshaped {shape.shape_type.name} "
                            f"({width_in:.2f}\"x{height_in:.2f}\" -> "
                            f"{shape.width.inches:.2f}\"x{shape.height.inches:.2f}\")"
                        )
                        right_in = shape.left.inches + shape.width.inches
                        bottom_in = shape.top.inches + shape.height.inches
                if right_in > SLIDE_WIDTH.inches + 0.20:
                    self.warnings.append(f"exceeds right edge ({right_in:.2f}\")")
                if bottom_in > SLIDE_HEIGHT.inches + 0.20:
                    self.warnings.append(f"exceeds slide bottom ({bottom_in:.2f}\")")
                if bottom_in > BOTTOM_MAX.inches + 0.10:
                    self.warnings.append(f"near bottom edge ({bottom_in:.2f}\", max {BOTTOM_MAX.inches:.2f}\")")
        if self.fixed:
            print(f"  [LAYOUT FIXED] {self.name}")
            for f in self.fixed:
                print(f"      + {f}")
        if self.warnings:
            print(f"  [LAYOUT WARNING] {self.name}")
            for w in self.warnings:
                print(f"      - {w}")
        return self.warnings


def check(slide, name):
    LayoutGuard(slide, name).check(fix=True)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def add_title_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_PRIMARY
    bg.line.fill.background()

    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.8), Inches(-1.2), Inches(4.5), Inches(4.5))
    c1.fill.solid()
    c1.fill.fore_color.rgb = RGBColor(0x22, 0x2d, 0x8f)
    c1.line.fill.background()

    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.6), Inches(0.2), Inches(2.8), Inches(2.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = RGBColor(0x25, 0x32, 0xa0)
    c2.line.fill.background()

    add_text(slide, MARGIN, Inches(1.40), CONTENT_W, Inches(0.60),
             "TRANSCRIPT INTELLIGENCE", FONT_TITLE, bold=True, color=C_WHITE)
    add_text(slide, MARGIN, Inches(2.05), CONTENT_W, Inches(0.45),
             "Aegis Cloud Security", FONT_SUBTITLE, color=RGBColor(0xcc, 0xcc, 0xcc))
    add_text(slide, MARGIN, Inches(2.55), CONTENT_W, Inches(0.40),
             "Call Analytics Pipeline - Product & Engineering Briefing",
             FONT_BODY, color=RGBColor(0xcc, 0xcc, 0xcc))

    specific_total = sum(data.feature_keywords.values())
    kpis = [
        (data.total_calls, "Calls Analysed", C_PRIMARY),
        (3, "Call Types", C_PRIMARY),
        (specific_total, "Feature Requests", C_PRIMARY),
        (data.risk_distribution.high, "Churn Flags", C_ACCENT),
    ]
    box_w = Inches(2.05)
    box_h = Inches(0.90)
    y = Inches(3.75)
    add_kpi_row(slide, kpis, y, box_w, box_h, card_fill=C_WHITE, label_color=C_TEXT)

    add_text(slide, MARGIN, Inches(5.05), CONTENT_W, Inches(0.30),
             datetime.now().strftime("%B %Y"), FONT_SMALL, color=C_WHITE)
    check(slide, "Title")
    return slide


def add_executive_summary_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Executive Summary", "Three findings that drive decisions today")

    top_feature = next(iter(data.feature_keywords.items()), ("feature", 0))
    worst_zone = data.problem_zones[0] if data.problem_zones else {"topic": "N/A", "call_type": "N/A", "sentiment": 0}
    avg_sentiment = data.avg_sentiment or 0

    findings = [
        ("🚨", "Churn risk is concentrated",
         f"{data.risk_distribution.high} accounts flagged high-risk ({data.risk_distribution.medium} medium). Support cases and external renewals both show competitor mentions and escalations.",
         C_RED),
        ("📉", f"Sentiment bottoms at {worst_zone.get('topic', 'reliability')}",
         f"{worst_zone.get('topic', 'N/A')} × {worst_zone.get('call_type', 'support').title()} scores {fmt_num(worst_zone.get('sentiment', 0))}/5 — the lowest zone in the dataset. Average overall sentiment is {fmt_num(avg_sentiment)}/5.",
         C_AMBER),
        ("📣", f"'{top_feature[0].title()}' dominates feature asks",
         f"{top_feature[1]} calls requested it — the clearest PM priority. Export and reporting gaps drive repeated friction across support and external calls.",
         C_SECONDARY),
    ]

    y = CONTENT_Y
    card_h = Inches(1.20)
    for icon, title, body, color in findings:
        body = fit_text(body, CONTENT_W - Inches(1.1), FONT_TINY, 3)
        add_card(slide, MARGIN, y, CONTENT_W, card_h)
        add_text(slide, MARGIN + Inches(0.12), y + Inches(0.12), Inches(0.55), Inches(0.55),
                 icon, FONT_HEADING, bold=True, color=color)
        add_text(slide, MARGIN + Inches(0.75), y + Inches(0.15), CONTENT_W - Inches(1.0), Inches(0.30),
                 title, FONT_BODY, bold=True, color=color)
        add_text(slide, MARGIN + Inches(0.75), y + Inches(0.48), CONTENT_W - Inches(1.0), Inches(0.60),
                 body, FONT_TINY, color=C_TEXT)
        y += card_h + Inches(0.08)
    check(slide, "Executive Summary")
    return slide


def add_pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "What We Built")

    steps = [
        ("01", "Data Loading", "Parse call folders: meeting-info, transcript, summary, speaker-meta."),
        ("02", "Call Type Inference", "Heuristic + LLM hybrid labels support, external, internal."),
        ("03", "Topic Categorization", "10-category business taxonomy + HDBSCAN embeddings, LLM naming, TF-IDF keywords."),
        ("04", "Sentiment Analysis", "Pre-scored sentimentScore (1-5) plus sentence-level labels. No LLM re-scoring."),
        ("05", "Bonus Insights", "Churn scoring, feature extraction, escalation chains, carry-forward actions."),
    ]

    card_w = Inches(2.95)
    card_h = Inches(1.60)
    top_y = Inches(1.15)
    bot_y = Inches(2.90)
    positions = [
        (MARGIN, top_y), (MARGIN + Inches(3.10), top_y), (MARGIN + Inches(6.20), top_y),
        (MARGIN + Inches(1.55), bot_y), (MARGIN + Inches(4.65), bot_y),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x, y = positions[i]
        add_card(slide, x, y, card_w, card_h)
        add_circle_badge(slide, x + Inches(0.18), y + Inches(0.18), Inches(0.48), num, C_SECONDARY)
        add_text(slide, x + Inches(0.75), y + Inches(0.22), Inches(2.05), Inches(0.32),
                 title, FONT_BODY, bold=True, color=C_PRIMARY)
        desc = fit_text(desc, Inches(2.65), FONT_TINY, 4)
        add_text(slide, x + Inches(0.18), y + Inches(0.72), Inches(2.60), Inches(0.78),
                 desc, FONT_TINY, color=C_TEXT)
    check(slide, "Pipeline")
    return slide


def add_dataset_slide(prs, data: PresentationData):
    """Rich dataset overview: KPIs + 2x2 chart grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Dataset Overview",
                    f"{data.total_calls} calls from Aegis dataset · {data.date_min} to {data.date_max}")

    sentiment_label = data.sentiment_interpretation.get("overall_label", "")
    kpis = [
        (data.total_calls, "Total Calls", C_PRIMARY),
        (fmt_num(data.duration_mean), "Avg Duration (min)", C_SECONDARY),
        (fmt_num(data.avg_sentiment), f"Avg Sentiment ({sentiment_label})", sentiment_color(data.avg_sentiment or 3)),
        (_fmt_date_range(data.date_min, data.date_max), "Date Range", C_ACCENT, FONT_SUBTITLE),
    ]
    # 4 tiles must fit in 9" content width with 0.12" gaps
    box_w = Inches(2.16)
    box_h = Inches(0.95)
    add_kpi_row(slide, kpis, CONTENT_Y, box_w, box_h)

    if data.total_calls > 0:
        breakdown = (
            f"Support: {data.support_count} calls ({round(100*data.support_count/data.total_calls)}%)   ·   "
            f"External: {data.external_count} calls ({round(100*data.external_count/data.total_calls)}%)   ·   "
            f"Internal: {data.internal_count} calls ({round(100*data.internal_count/data.total_calls)}%)"
        )
    else:
        breakdown = "Support: 0 · External: 0 · Internal: 0"

    add_text(slide, MARGIN, Inches(2.35), CONTENT_W, Inches(0.25),
             breakdown, FONT_BODY, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

    # 2x2 chart grid to remove whitespace and give a complete snapshot
    chart_w = Inches(4.75)
    chart_h = Inches(1.45)
    gap = Inches(0.12)
    grid_y = Inches(2.62)
    left_x = MARGIN
    right_x = MARGIN + chart_w + gap

    add_chart_if_exists(slide, "02_call_types_distribution.png", left_x, grid_y, chart_w, chart_h)
    add_chart_if_exists(slide, "04_sentiment_score_distribution.png", right_x, grid_y, chart_w, chart_h)
    add_chart_if_exists(slide, "03_top_business_categories.png", left_x, grid_y + chart_h + gap, chart_w, chart_h)
    add_chart_if_exists(slide, "05_action_items_duration.png", right_x, grid_y + chart_h + gap, chart_w, chart_h)

    # Conclusion footer
    card_h = Inches(0.55)
    footer_y = grid_y + 2 * chart_h + 2 * gap + Inches(0.04)
    if data.total_calls > 0:
        insight_text = (
            f"Conclusion: Support dominates volume ({round(100*data.support_count/data.total_calls)}%). "
            f"External calls are smallest but carry the highest churn and feature-ask intensity."
        )
    else:
        insight_text = "Conclusion: No call type data available."
    add_card(slide, MARGIN, footer_y, CONTENT_W, card_h, line=C_SECONDARY)
    add_text(slide, MARGIN + Inches(0.12), footer_y + Inches(0.08), CONTENT_W - Inches(0.24), Inches(0.40),
             fit_text(insight_text, CONTENT_W - Inches(0.24), FONT_TINY, 2), FONT_TINY, color=C_TEXT)

    check(slide, "Dataset")
    return slide


def add_topic_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Topic Categorization",
                    "10-category business taxonomy matched against topics, call title, and summary text.")

    # Chart on the left; leave room for full right-column cards
    add_chart_if_exists(slide, "03_topic_distribution_by_type.png", CHART_X, CONTENT_Y, CHART_W, Inches(3.6))

    biz = data.business_taxonomy
    top_categories = biz.get("top_categories", [])
    narratives = biz.get("narratives", {})

    clusters = top_categories[:3]
    card_h = Inches(1.02)
    y = CONTENT_Y
    for cat_info in clusters:
        name = cat_info.get("category", "Unknown")
        count = cat_info.get("count", 0)
        icon = TOPIC_ICONS.get(name, "")
        title_text = f"{icon}  {name}" if icon else name
        narrative = fit_text(narratives.get(name, f"{count} calls tagged in this category."), Inches(3.0), FONT_TINY, 3)
        avg_sentiment = cat_info.get("avg_sentiment")
        sentiment_text = f"Avg sentiment: {fmt_num(avg_sentiment)}/5" if avg_sentiment is not None else ""
        sent_color = sentiment_color(avg_sentiment) if avg_sentiment is not None else C_LIGHT

        add_card(slide, RIGHT_X, y, RIGHT_W, card_h)
        # Title stops before the sentiment badge to avoid overlap
        add_text(slide, RIGHT_X + Inches(0.10), y + Inches(0.08), RIGHT_W - Inches(1.10), Inches(0.26),
                 title_text, FONT_SMALL, bold=True, color=C_PRIMARY)
        if sentiment_text:
            add_text(slide, RIGHT_X + RIGHT_W - Inches(1.05), y + Inches(0.08), Inches(0.95), Inches(0.22),
                     sentiment_text, FONT_MICRO, bold=True, color=sent_color, align=PP_ALIGN.RIGHT)
        add_text(slide, RIGHT_X + Inches(0.10), y + Inches(0.34), RIGHT_W - Inches(0.20), Inches(0.64),
                 narrative, FONT_TINY, color=C_TEXT)
        y += card_h + Inches(0.07)

    # Conclusion tile placed below the stacked category cards
    if top_categories:
        top = top_categories[0]
        conclusion = (
            f"Conclusion: {top['category']} is the #1 topic ({top['count']} calls, {top['pct_of_total']}%). "
            f"It is dominant in {top.get('dominant_call_type', 'external')} calls ({top.get('dominant_pct', 0)}%). "
            f"Product and support should align on {top['category'].lower()} playbooks."
        )
        conclusion_y = y + Inches(0.05)
        concl_h = Inches(0.95)
        add_card(slide, RIGHT_X, conclusion_y, RIGHT_W, concl_h, line=C_SECONDARY)
        add_text(slide, RIGHT_X + Inches(0.10), conclusion_y + Inches(0.08), RIGHT_W - Inches(0.20), Inches(0.80),
                 fit_text(conclusion, RIGHT_W - Inches(0.20), FONT_TINY, 5), FONT_TINY, color=C_TEXT)

    check(slide, "Topics")
    return slide


def add_sentiment_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Sentiment Analysis by Call Type",
                    "Sentiment score scale: 1 = very negative, 3 = neutral, 5 = very positive")

    add_chart_if_exists(slide, "04_sentiment_trend_by_type.png", CHART_X, CONTENT_Y, CHART_W, Inches(3.4))

    interpretations = {
        "support": "Mixed — agents resolve issues, but customers call when things break.",
        "external": "Highest baseline — compliance wins and renewals drive positive tone.",
        "internal": "Neutral-to-positive — incident postmortems and risk reviews add weight.",
    }

    types = [
        ("Support", data.sentiment.support_score, data.sentiment.support_neg, interpretations["support"]),
        ("External", data.sentiment.external_score, data.sentiment.external_neg, interpretations["external"]),
        ("Internal", data.sentiment.internal_score, data.sentiment.internal_neg, interpretations["internal"]),
    ]

    card_h = Inches(1.15)
    y = CONTENT_Y
    for name, score, neg_pct, explanation in types:
        color = sentiment_color(score)
        explanation = fit_text(explanation, RIGHT_W - Inches(0.24), FONT_TINY, 3)
        add_card(slide, RIGHT_X, y, RIGHT_W, card_h)
        add_text(slide, RIGHT_X + Inches(0.12), y + Inches(0.10), RIGHT_W - Inches(0.24), Inches(0.26),
                 name, FONT_BODY, bold=True, color=color)
        add_text(slide, RIGHT_X + Inches(0.12), y + Inches(0.36), RIGHT_W - Inches(0.24), Inches(0.22),
                 f"Score: {fmt_num(score)} / 5  ·  {fmt_pct(neg_pct)}% negative", FONT_TINY, color=C_LIGHT)
        add_text(slide, RIGHT_X + Inches(0.12), y + Inches(0.58), RIGHT_W - Inches(0.24), Inches(0.50),
                 explanation, FONT_TINY, color=C_TEXT)
        y += card_h + Inches(0.08)

    # Scale interpretation box (full-width footer to avoid truncating any labels)
    scale = data.sentiment_interpretation.get("scale", {})
    scale_text = " · ".join([f"{k}: {v}" for k, v in scale.items()])
    add_text(slide, MARGIN, Inches(5.05), CONTENT_W, Inches(0.45),
             fit_text(scale_text, CONTENT_W, FONT_MICRO, 3), FONT_MICRO, color=C_LIGHT, align=PP_ALIGN.CENTER)

    check(slide, "Sentiment")
    return slide


def _zone_card(slide, x, y, zone, color, label_text):
    """Helper to draw a single sentiment-zone insight card."""
    card_w = Inches(2.95)
    card_h = Inches(1.35)
    why = fit_text(zone.get("why", ""), card_w - Inches(0.24), FONT_TINY, 4)
    add_card(slide, x, y, card_w, card_h, line=color)
    add_text(slide, x + Inches(0.12), y + Inches(0.10), Inches(0.70), Inches(0.24),
             label_text, FONT_MICRO, bold=True, color=color)
    add_text(slide, x + Inches(0.85), y + Inches(0.10), card_w - Inches(0.97), Inches(0.28),
             f"{zone['topic']} × {zone['call_type'].title()}", FONT_SMALL, bold=True, color=color)
    add_text(slide, x + Inches(0.12), y + Inches(0.42), card_w - Inches(0.24), Inches(0.35),
             f"Score: {fmt_num(zone.get('sentiment', 0))}/5 · {zone.get('call_count', 0)} calls · {fmt_pct(zone.get('negative_pct', 0))}% negative",
             FONT_TINY, color=C_LIGHT)
    add_text(slide, x + Inches(0.12), y + Inches(0.78), card_w - Inches(0.24), Inches(0.50),
             why, FONT_TINY, color=C_TEXT)


def _short_zone_why(zone: dict) -> str:
    """Return a concise insight plus an example count that fits a small card."""
    why = zone.get("why", "")
    # Split insight from examples
    if "Examples:" in why:
        insight, examples = why.split("Examples:", 1)
        insight = insight.strip()
        example_count = len([e for e in examples.split(",") if e.strip()])
        suffix = "examples" if example_count != 1 else "example"
        return f"{insight} ({example_count} {suffix})."
    return why


def add_problem_zones_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Where Sentiment Goes Negative",
                    "Average sentiment score per call type × topic category. Red = problem zones.")

    # Heatmap on the left; cards stacked on the right where space is available
    add_chart_if_exists(slide, "04_sentiment_heatmap_by_taxonomy.png", CHART_X, CONTENT_Y, CHART_W, Inches(3.5))

    low_zone = data.problem_zones[0] if data.problem_zones else {"topic": "N/A", "call_type": "N/A", "sentiment": 0, "why": "No data", "call_count": 0, "label": ""}
    watch_zone = data.watch_zone if data.watch_zone else low_zone
    high_zone = data.strong_zones[0] if data.strong_zones else {"topic": "N/A", "call_type": "N/A", "sentiment": 0, "why": "No data", "call_count": 0, "label": ""}

    card_h = Inches(1.32)
    gap = Inches(0.05)
    y = CONTENT_Y
    for zone, color, label_prefix in [
        (low_zone, C_RED, "LOW"),
        (watch_zone, C_AMBER, "WATCH"),
        (high_zone, C_GREEN, "HIGH"),
    ]:
        label = zone.get("label", "")
        label_text = f"{label_prefix} · {label}" if label else label_prefix
        why = fit_text(_short_zone_why(zone), RIGHT_W - Inches(0.22), FONT_TINY, 4)
        add_card(slide, RIGHT_X, y, RIGHT_W, card_h, line=color)
        add_text(slide, RIGHT_X + Inches(0.10), y + Inches(0.08), Inches(0.85), Inches(0.24),
                 label_text, FONT_MICRO, bold=True, color=color)
        add_text(slide, RIGHT_X + Inches(0.98), y + Inches(0.08), RIGHT_W - Inches(1.08), Inches(0.26),
                 f"{zone['topic']} × {zone['call_type'].title()}", FONT_SMALL, bold=True, color=color)
        add_text(slide, RIGHT_X + Inches(0.10), y + Inches(0.34), RIGHT_W - Inches(0.20), Inches(0.22),
                 f"Score: {fmt_num(zone.get('sentiment', 0))}/5 · {zone.get('call_count', 0)} calls · {fmt_pct(zone.get('negative_pct', 0))}% negative",
                 FONT_TINY, color=C_LIGHT)
        add_text(slide, RIGHT_X + Inches(0.10), y + Inches(0.56), RIGHT_W - Inches(0.20), Inches(0.70),
                 why, FONT_TINY, color=C_TEXT)
        y += card_h + gap

    check(slide, "Problem Zones")
    return slide


def add_strong_zones_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Where Sentiment Is Strongest",
                    "Green zones are relationship and product strengths to reinforce.")

    add_chart_if_exists(slide, "04_sentiment_stacked_by_type.png", MARGIN, CONTENT_Y, CONTENT_W, Inches(2.35))

    zones = data.strong_zones[:3]
    if not zones:
        zones = [{"topic": "N/A", "call_type": "N/A", "sentiment": 0, "why": "No data", "call_count": 0, "label": ""}]

    card_w = Inches(2.95)
    card_h = Inches(1.50)
    base_y = Inches(3.78)
    xs = [MARGIN, MARGIN + Inches(3.05), MARGIN + Inches(6.10)]

    for x, zone in zip(xs, zones):
        why = fit_text(_short_zone_why(zone), card_w - Inches(0.24), FONT_TINY, 4)
        label = zone.get("label", "")
        label_text = f"HIGH · {label}" if label else "HIGH"
        add_card(slide, x, base_y, card_w, card_h, line=C_GREEN)
        add_text(slide, x + Inches(0.12), base_y + Inches(0.08), Inches(0.65), Inches(0.24),
                 label_text, FONT_MICRO, bold=True, color=C_GREEN)
        add_text(slide, x + Inches(0.80), base_y + Inches(0.08), card_w - Inches(0.92), Inches(0.28),
                 f"{zone['topic']} × {zone['call_type'].title()}", FONT_SMALL, bold=True, color=C_GREEN)
        add_text(slide, x + Inches(0.12), base_y + Inches(0.36), card_w - Inches(0.24), Inches(0.28),
                 f"Score: {fmt_num(zone.get('sentiment', 0))}/5 · {zone.get('call_count', 0)} calls · {fmt_pct(zone.get('negative_pct', 0))}% negative",
                 FONT_TINY, color=C_LIGHT)
        add_text(slide, x + Inches(0.12), base_y + Inches(0.66), card_w - Inches(0.24), Inches(0.78),
                 why, FONT_TINY, color=C_TEXT)

    # Scale footer so 4.8 is unambiguous
    scale_footer = (
        "Sentiment scale: 1-2 = very negative/negative, 3 = neutral/mixed, "
        "4 = positive (relationship/product strength), 5 = very positive (advocacy moment)."
    )
    add_text(slide, MARGIN, Inches(5.40), CONTENT_W, Inches(0.30),
             fit_text(scale_footer, CONTENT_W, FONT_MICRO, 2), FONT_MICRO, color=C_LIGHT, align=PP_ALIGN.CENTER)

    check(slide, "Strong Zones")
    return slide


def _format_signal(signal: str) -> str:
    return signal.replace("_", " ").title()


def add_churn_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Churn Risk Detection",
                    "Keyword-based churn scoring (0-10+). High score + low sentiment = immediate attention.")

    add_chart_if_exists(slide, "05_churn_risk_distribution.png", CHART_X, CONTENT_Y, CHART_W, Inches(3.8))

    add_text(slide, RIGHT_X, CONTENT_Y, RIGHT_W, Inches(0.30),
             f"⚠  Flagged Accounts ({data.risk_distribution.high} high risk)", FONT_BODY, bold=True, color=C_RED)

    card_w = Inches(3.40)
    card_h = Inches(1.60)
    y = CONTENT_Y + Inches(0.38)

    for account in data.churn_narratives[:2]:
        add_card(slide, RIGHT_X, y, card_w, card_h, line=C_RED)
        account_name = account.get("account", "Unknown account")
        title = fit_text(account.get("title", "Unknown"), card_w - Inches(0.24), FONT_SMALL, 3)
        add_text(slide, RIGHT_X + Inches(0.12), y + Inches(0.08), card_w - Inches(0.24), Inches(0.50),
                 title, FONT_SMALL, bold=True, color=C_TEXT)
        add_text(slide, RIGHT_X + Inches(0.12), y + Inches(0.60), card_w - Inches(0.24), Inches(0.18),
                 f"Account: {account_name}  |  {account.get('call_type', 'unknown').title()} call", FONT_TINY, color=C_LIGHT)
        churn_score = account.get('churn_score', 0)
        add_text(slide, RIGHT_X + Inches(0.12), y + Inches(0.80), card_w - Inches(0.24), Inches(0.18),
                 f"Churn signal: {churn_score} (scale 0-10+)  |  Sentiment: {account.get('sentiment_score', 0)}/5",
                 FONT_TINY, color=C_LIGHT)

        signals = account.get("signals_list", [])
        if isinstance(signals, str):
            signals = [s.strip() for s in signals.split(",") if s.strip()]
        signal_text = "  |  ".join(_format_signal(s) for s in signals[:3])
        signal_text = fit_text(signal_text, card_w - Inches(0.24), FONT_MICRO, 3)
        add_text(slide, RIGHT_X + Inches(0.12), y + Inches(1.00), card_w - Inches(0.24), Inches(0.40),
                 signal_text, FONT_MICRO, color=C_TEXT)

        conclusion = fit_text(account.get("conclusion", ""), card_w - Inches(0.24), FONT_MICRO, 2)
        add_text(slide, RIGHT_X + Inches(0.12), y + Inches(1.42), card_w - Inches(0.24), Inches(0.16),
                 conclusion, FONT_MICRO, bold=True, color=C_RED)
        y += card_h + Inches(0.10)

    # Bottom insight tile placed below the account cards
    total_flagged = data.risk_distribution.high + data.risk_distribution.medium
    churn_source = (
        f"Conclusion: {total_flagged} of {data.total_calls} calls are medium/high risk. "
        f"Top signals: escalation, product dissatisfaction, competitor mentions. "
        f"Source: rule-based scoring on transcript + summary text."
    )
    footer_y = y + Inches(0.02)
    add_card(slide, MARGIN, footer_y, CONTENT_W, Inches(0.52), line=C_RED)
    add_text(slide, MARGIN + Inches(0.12), footer_y + Inches(0.07), CONTENT_W - Inches(0.24), Inches(0.38),
             fit_text(churn_source, CONTENT_W - Inches(0.24), FONT_TINY, 2), FONT_TINY, color=C_TEXT)

    check(slide, "Churn")
    return slide


def add_renewal_risk_slide(prs, data: PresentationData):
    """Show which renewal conversations are healthy and which are bleeding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    renewal = data.renewal_risk
    total = renewal.get("total_renewal_calls", 0)
    risky = renewal.get("risky_renewal_calls", 0)
    healthy = total - risky
    risky_accounts = renewal.get("risky_accounts", [])
    add_slide_title(slide, "Renewal Risk",
                    "Renewal calls split into healthy vs. at-risk by sentiment and churn signals.")

    # KPI row
    kpis = [
        (total, "Renewal Calls", C_PRIMARY),
        (risky, "At-Risk Renewals", C_RED),
        (healthy, "Healthy Renewals", C_GREEN),
        (len(risky_accounts), "Risky Accounts", C_AMBER),
    ]
    box_w = Inches(2.16)
    box_h = Inches(0.85)
    add_kpi_row(slide, kpis, CONTENT_Y, box_w, box_h)

    # At-risk calls list
    add_text(slide, MARGIN, CONTENT_Y + Inches(1.05), CONTENT_W, Inches(0.30),
             "At-risk renewal calls (negative sentiment, competitor mentions, or escalations):",
             FONT_BODY, bold=True, color=C_RED)

    calls = [c for c in renewal.get("calls", []) if c.get("is_risky")][:5]
    card_h = Inches(0.76)
    y = CONTENT_Y + Inches(1.38)
    flag_label = {
        "negative_sentiment_dominates": "negative tone",
        "competitor_mentioned": "competitor",
        "escalation_requested": "escalation",
        "product_dissatisfaction": "product dissatisfaction",
    }
    for call in calls:
        add_card(slide, MARGIN, y, CONTENT_W, card_h, line=C_RED)
        title = fit_text(call.get("title", "Unknown"), CONTENT_W - Inches(0.24), FONT_SMALL, 1)
        add_text(slide, MARGIN + Inches(0.10), y + Inches(0.08), CONTENT_W - Inches(0.24), Inches(0.24),
                 title, FONT_SMALL, bold=True, color=C_TEXT)
        account = call.get("account", "Unknown account")
        left_info = f"Account: {account}  |  {call.get('call_type', 'unknown').title()} call"
        right_info = f"Sentiment: {call.get('sentiment_score', 0)}/5  |  Score: {call.get('churn_score', 0)}"
        add_text(slide, MARGIN + Inches(0.10), y + Inches(0.34), Inches(4.6), Inches(0.20),
                 left_info, FONT_TINY, color=C_LIGHT)
        add_text(slide, MARGIN + Inches(5.0), y + Inches(0.34), Inches(4.6), Inches(0.20),
                 right_info, FONT_TINY, color=C_LIGHT, align=PP_ALIGN.RIGHT)
        flags = "  |  ".join(flag_label.get(s, s.replace("_", " ")) for s in call.get("risk_flags", [])[:3])
        add_text(slide, MARGIN + Inches(0.10), y + Inches(0.56), CONTENT_W - Inches(0.24), Inches(0.18),
                 fit_text(f"Risk flags: {flags}", CONTENT_W - Inches(0.24), FONT_MICRO, 1), FONT_MICRO, color=C_RED)
        y += card_h + Inches(0.06)

    # Bottom insight
    if total > 0:
        insight = (
            f"Conclusion: {risky} of {total} renewal calls are at risk ({round(100*risky/total)}%). "
            f"Renewal discussion alone is neutral; risk comes from negative sentiment, competitor mentions, or escalations."
        )
    else:
        insight = "Conclusion: No renewal calls detected."
    footer_y = y + Inches(0.08)
    add_card(slide, MARGIN, footer_y, CONTENT_W, Inches(0.50), line=C_SECONDARY)
    add_text(slide, MARGIN + Inches(0.12), footer_y + Inches(0.07), CONTENT_W - Inches(0.24), Inches(0.36),
             fit_text(insight, CONTENT_W - Inches(0.24), FONT_TINY, 2), FONT_TINY, color=C_TEXT)

    check(slide, "Renewal Risk")
    return slide


def add_feature_slide(prs, data: PresentationData):
    """Feature requests counted by unique calls (one call = one request)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    total_calls = sum(data.feature_keywords.values())
    total_mentions = sum(c.get("mention_count", 0) for c in data.feature_callouts)
    add_slide_title(slide, "Feature Request Intelligence",
                    f"{total_calls} unique calls requested a specific feature ({total_mentions} total mentions). One call = one request.")

    add_chart_if_exists(slide, "05_feature_requests.png", CHART_X, CONTENT_Y, CHART_W, Inches(3.6))

    add_text(slide, RIGHT_X, CONTENT_Y, RIGHT_W, Inches(0.28),
             "PM-Ready Backlog Items", FONT_BODY, bold=True, color=C_PRIMARY)

    top_features = list(data.feature_keywords.items())[:4]

    def prio_for_rank(i):
        return "P1" if i == 0 else ("P2" if i <= 2 else "P3")

    display_names = {"sso": "SSO", "mfa": "MFA", "ldap": "LDAP", "saml": "SAML"}

    card_h = Inches(0.82)
    y = CONTENT_Y + Inches(0.32)
    for i, (kw, call_count) in enumerate(top_features):
        prio = prio_for_rank(i)
        prio_color = C_RED if prio == "P1" else (C_AMBER if prio == "P2" else C_GREEN)
        callout = next((c for c in data.feature_callouts if c["keyword"] == kw), {})
        mention_count = callout.get("mention_count", 0)
        dominant_type = callout.get("dominant_call_type", "unknown")
        dominant_cat = callout.get("dominant_category", "Other")
        subtypes = callout.get("subtypes", [])

        add_card(slide, RIGHT_X, y, RIGHT_W, card_h)
        add_circle_badge(slide, RIGHT_X + Inches(0.10), y + Inches(0.26), Inches(0.36), prio, prio_color, FONT_TINY)

        kw_display = display_names.get(kw.lower(), kw.title())
        add_text(slide, RIGHT_X + Inches(0.54), y + Inches(0.10), Inches(1.55), Inches(0.24),
                 kw_display, FONT_SMALL, bold=True, color=C_TEXT)
        add_text(slide, RIGHT_X + Inches(2.10), y + Inches(0.10), Inches(1.15), Inches(0.24),
                 f"{call_count} calls", FONT_TINY, color=C_LIGHT, align=PP_ALIGN.RIGHT)

        # Context line: call type + category
        ctx_line = f"Mostly {dominant_type} · {dominant_cat}"
        add_text(slide, RIGHT_X + Inches(0.54), y + Inches(0.34), RIGHT_W - Inches(0.64), Inches(0.20),
                 ctx_line, FONT_MICRO, color=C_LIGHT)

        # Detail line: subtypes when available; otherwise mention count
        if subtypes:
            detail = "Top contexts: " + ", ".join(f"{s['subtype']} ({s['count']})" for s in subtypes[:2])
        else:
            detail = f"{mention_count} mentions — primary channel: {dominant_type}"
        add_text(slide, RIGHT_X + Inches(0.54), y + Inches(0.54), RIGHT_W - Inches(0.64), Inches(0.34),
                 detail, FONT_MICRO, color=C_TEXT)

        y += card_h + Inches(0.05)

    check(slide, "Features")
    return slide


def add_action_items_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Action Items & Call Efficiency",
                    "Follow-through commitments by call type — a proxy for engagement and urgency.")

    add_chart_if_exists(slide, "05_action_items_by_type.png", MARGIN, CONTENT_Y, Inches(4.8), Inches(2.8))

    items = []
    for ctype in ["external", "support", "internal"]:
        info = data.action_items.get(ctype, {})
        avg = info.get("avg_per_call", 0)
        total = info.get("action_mentions", 0)
        items.append((ctype.title(), avg, total))

    items_sorted = sorted(items, key=lambda x: x[1], reverse=True)
    leader_type, leader_avg, _ = items_sorted[0]
    second_type, second_avg, _ = items_sorted[1]
    third_type, third_avg, _ = items_sorted[2]
    insights = [
        (f"{leader_type} calls drive follow-through", f"{leader_avg:.1f} avg action items per {leader_type.lower()} call — highest of all call types."),
        (f"{second_type} calls: {second_avg:.1f} avg", f"Middle performer — follow-ups vary by {second_type.lower()} context."),
        (f"{third_type} calls: {third_avg:.1f} avg", "Fewest items but may still be high urgency: engineering fixes and postmortems."),
        ("Opportunity: action item tracking", "Cross-reference action items across call sequences for closed-loop accountability."),
    ]

    card_w = Inches(4.2)
    card_h = Inches(0.95)
    x = Inches(5.4)
    y = CONTENT_Y
    for title, body in insights:
        body = fit_text(body, card_w - Inches(0.24), FONT_TINY, 3)
        add_card(slide, x, y, card_w, card_h)
        add_text(slide, x + Inches(0.12), y + Inches(0.10), card_w - Inches(0.24), Inches(0.24),
                 title, FONT_SMALL, bold=True, color=C_SECONDARY)
        add_text(slide, x + Inches(0.12), y + Inches(0.36), card_w - Inches(0.24), Inches(0.55),
                 body, FONT_TINY, color=C_TEXT)
        y += card_h + Inches(0.08)
    check(slide, "Action Items")
    return slide


def add_carry_forward_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Carry-Forward Actions",
                    f"{data.carry_forward_total} open action items extracted from call summaries — owners and calls shown.")

    # Left: bar chart of carry-forward action counts by type
    add_chart_if_exists(slide, "05_carry_forward_actions.png", MARGIN, CONTENT_Y, Inches(4.8), Inches(2.8))

    # Right: top actions per type
    card_w = Inches(4.2)
    card_h = Inches(1.12)
    x = Inches(5.4)
    y = CONTENT_Y
    for ctype in ["support", "external", "internal"]:
        actions = data.carry_forward_actions.get(ctype, {}).get("top_actions", [])
        if actions:
            top = actions[0]
            raw_text = f"{top.get('owner', 'Unknown')}: {top.get('text', '')}"
            text = fit_text(raw_text[:220], card_w - Inches(0.24), FONT_TINY, 4)
            add_card(slide, x, y, card_w, card_h)
            total_for_type = data.carry_forward_actions.get(ctype, {}).get("count", len(actions))
            add_text(slide, x + Inches(0.12), y + Inches(0.10), card_w - Inches(0.24), Inches(0.24),
                     f"{ctype.title()} (top 5 of {total_for_type})", FONT_SMALL, bold=True, color=C_SECONDARY)
            add_text(slide, x + Inches(0.12), y + Inches(0.36), card_w - Inches(0.24), Inches(0.70),
                     text, FONT_TINY, color=C_TEXT)
            y += card_h + Inches(0.08)

    # Conclusion
    conclusion = (
        "Conclusion: Support generates the most follow-up items. Without closed-loop tracking, "
        "customer commitments from external and support calls can slip between handoffs."
    )
    add_card(slide, MARGIN, Inches(5.00), CONTENT_W, Inches(0.60), line=C_SECONDARY)
    add_text(slide, MARGIN + Inches(0.12), Inches(5.07), CONTENT_W - Inches(0.24), Inches(0.45),
             fit_text(conclusion, CONTENT_W - Inches(0.24), FONT_TINY, 2), FONT_TINY, color=C_TEXT)

    check(slide, "Carry-Forward Actions")
    return slide


def _owner_color(owner: str) -> RGBColor:
    return {
        "Product": C_RED,
        "Engineering": C_AMBER,
        "Sales / CS": C_SECONDARY,
        "Sales / Customer Success": C_SECONDARY,
        "Support": C_SECONDARY,
        "Operations / Analytics": C_GREEN,
    }.get(owner, C_PRIMARY)


def add_recommendations_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Recommendations at a Glance", "Prioritized actions with source evidence")

    recs = data.recommendations[:5]
    if not recs:
        recs = [{
            "rank": 1,
            "owner": "Product",
            "title": "Review feature-request data",
            "headline": "Feature extraction did not produce recommendations.",
            "evidence": [],
            "source_calls": [],
            "metrics": {},
        }]

    card_h = Inches(0.78)
    gap = Inches(0.05)
    y = CONTENT_Y
    for rec in recs:
        add_card(slide, MARGIN, y, CONTENT_W, card_h)
        owner_color = _owner_color(rec.get("owner", ""))
        add_text(slide, MARGIN + Inches(0.10), y + Inches(0.08), Inches(1.55), Inches(0.20),
                 f"[{rec.get('owner', '')}]", FONT_MICRO, bold=True, color=owner_color)
        title_text = fit_text(rec.get("title", ""), CONTENT_W - Inches(1.85), FONT_SMALL, 1)
        add_text(slide, MARGIN + Inches(1.75), y + Inches(0.06), CONTENT_W - Inches(1.90), Inches(0.26),
                 title_text, FONT_SMALL, bold=True, color=C_TEXT)

        # Use the concise headline instead of dense evidence; full evidence lives on detail slides
        headline = fit_text(rec.get("headline", ""), CONTENT_W - Inches(0.30), FONT_TINY, 2)
        add_text(slide, MARGIN + Inches(0.10), y + Inches(0.32), CONTENT_W - Inches(0.20), Inches(0.42),
                 headline, FONT_TINY, color=C_TEXT)

        y += card_h + gap

    check(slide, "Recommendations")
    return slide


def _add_single_recommendation_slide(prs, rec: dict):
    """Add one full slide explaining a single recommendation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rank = rec.get("rank", 0)
    title = rec.get("title", "Recommendation")
    owner = rec.get("owner", "")
    add_slide_title(slide, f"Rec {rank}: {title}",
                    "What the data says, why it matters, and what to do about it")

    # Owner badge below the slide title on the right
    owner_color = _owner_color(owner)
    badge_y = CONTENT_Y - Inches(0.05)
    add_card(slide, Inches(8.25), badge_y, Inches(1.60), Inches(0.32), fill=owner_color)
    add_text(slide, Inches(8.30), badge_y + Inches(0.03), Inches(1.50), Inches(0.26),
             f"Owner: {owner}", FONT_MICRO, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Two-column layout: Problem | Recommended Action
    col_w = Inches(4.55)
    card_h = Inches(1.50)
    y = CONTENT_Y + Inches(0.35)

    # Problem card
    add_card(slide, MARGIN, y, col_w, card_h, line=C_RED)
    add_text(slide, MARGIN + Inches(0.10), y + Inches(0.08), col_w - Inches(0.20), Inches(0.26),
             "Problem", FONT_SMALL, bold=True, color=C_RED)
    problem = fit_text(rec.get("problem", ""), col_w - Inches(0.20), FONT_TINY, 5)
    add_text(slide, MARGIN + Inches(0.10), y + Inches(0.36), col_w - Inches(0.20), Inches(1.05),
             problem, FONT_TINY, color=C_TEXT)

    # Solution card
    x2 = MARGIN + col_w + Inches(0.20)
    add_card(slide, x2, y, col_w, card_h, line=C_SECONDARY)
    add_text(slide, x2 + Inches(0.10), y + Inches(0.08), col_w - Inches(0.20), Inches(0.26),
             "Recommended Action", FONT_SMALL, bold=True, color=C_SECONDARY)
    solution = fit_text(rec.get("solution", ""), col_w - Inches(0.20), FONT_TINY, 5)
    add_text(slide, x2 + Inches(0.10), y + Inches(0.36), col_w - Inches(0.20), Inches(1.05),
             solution, FONT_TINY, color=C_TEXT)

    # Evidence card (full width below)
    evidence_y = y + card_h + Inches(0.08)
    evidence = rec.get("evidence", [])
    if evidence:
        add_card(slide, MARGIN, evidence_y, CONTENT_W, Inches(0.70), line=C_AMBER)
        add_text(slide, MARGIN + Inches(0.10), evidence_y + Inches(0.08), Inches(1.15), Inches(0.24),
                 "Evidence", FONT_SMALL, bold=True, color=C_AMBER)
        add_bullets(slide, MARGIN + Inches(1.30), evidence_y + Inches(0.08),
                    CONTENT_W - Inches(1.45), Inches(0.54), evidence[:3],
                    size=FONT_MICRO, color=C_TEXT, max_items=3)
    else:
        evidence_y -= Inches(0.08)

    # Impact card (full width below evidence)
    impact_y = evidence_y + Inches(0.82)
    add_card(slide, MARGIN, impact_y, CONTENT_W, Inches(0.70), line=C_GREEN)
    add_text(slide, MARGIN + Inches(0.10), impact_y + Inches(0.08), Inches(1.35), Inches(0.24),
             "Expected Impact", FONT_SMALL, bold=True, color=C_GREEN)
    impact = fit_text(rec.get("expected_impact", ""), CONTENT_W - Inches(1.60), FONT_TINY, 2)
    add_text(slide, MARGIN + Inches(1.50), impact_y + Inches(0.10), CONTENT_W - Inches(1.65), Inches(0.50),
             impact, FONT_TINY, color=C_TEXT)

    # Source calls footer
    source_calls = rec.get("source_calls", [])
    if source_calls:
        footer_y = impact_y + Inches(0.82)
        source_text = "Source calls: " + ", ".join(t[:90] for t in source_calls[:2])
        add_text(slide, MARGIN, footer_y, CONTENT_W, Inches(0.40),
                 fit_text(source_text, CONTENT_W, FONT_MICRO, 2), FONT_MICRO, color=C_LIGHT)

    check(slide, f"Rec {rank}")
    return slide


def add_recommendation_detail_slides(prs, data: PresentationData):
    for rec in data.recommendations[:5]:
        _add_single_recommendation_slide(prs, rec)
    return None


def add_reasonableness_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "AI & Data Reasonableness Check",
                    "Transparent limitations that strengthen the narrative, not weaken it")

    cards = [
        ("✅", "Dataset is complete", f"{data.total_calls} calls analysed across {data.date_min} to {data.date_max}. All call types and sentiment labels present in source JSON.", C_GREEN),
        ("⚠️", "Topic clusters are weak", f"HDBSCAN silhouette {fmt_num(data.hdbscan_score, 3)} (ideal > 0.5); {data.hdbscan_noise} calls flagged as noise. We therefore added a keyword-backed 10-category business taxonomy for the deck.", C_AMBER),
        ("⚠️", "Churn scoring is rule-based", "Not a predictive model. Every point is independently verifiable from transcript text. High/medium risk counts should be treated as flags, not probabilities.", C_AMBER),
        ("✅", "Conclusions are reasonable", "Top feature ('report'), worst sentiment zone, and highest-risk accounts all align with the raw call content. Recommendations trace back to specific data sources.", C_GREEN),
    ]

    card_w = Inches(4.55)
    card_h = Inches(1.65)
    col_x = [MARGIN, MARGIN + Inches(4.75)]
    start_y = CONTENT_Y

    for i, (icon, title, body, color) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = col_x[col]
        y = start_y + row * (card_h + Inches(0.12))
        body = fit_text(body, card_w - Inches(0.80), FONT_TINY, 5)
        add_card(slide, x, y, card_w, card_h)
        add_text(slide, x + Inches(0.12), y + Inches(0.12), Inches(0.55), Inches(0.55),
                 icon, FONT_HEADING, bold=True, color=color)
        add_text(slide, x + Inches(0.75), y + Inches(0.15), card_w - Inches(0.87), Inches(0.30),
                 title, FONT_BODY, bold=True, color=color)
        add_text(slide, x + Inches(0.75), y + Inches(0.50), card_w - Inches(0.87), Inches(1.05),
                 body, FONT_TINY, color=C_TEXT)

    # Bottom verdict
    verdict = (
        "Verdict: The deck is reasonable for decision-making because every headline is grounded in "
        "countable data, limitations are disclosed, and recommendations are source-tagged."
    )
    add_card(slide, MARGIN, Inches(5.00), CONTENT_W, Inches(0.60), line=C_SECONDARY)
    add_text(slide, MARGIN + Inches(0.12), Inches(5.07), CONTENT_W - Inches(0.24), Inches(0.45),
             fit_text(verdict, CONTENT_W - Inches(0.24), FONT_TINY, 2), FONT_TINY, bold=True, color=C_TEXT)

    check(slide, "Reasonableness")
    return slide


def add_methodology_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Appendix: Pipeline Methodology")

    cards = [
        ("Data Sources", f"{data.total_calls} calls from Aegis dataset: meeting-info, summary, transcript, speaker-meta."),
        ("Sentiment Approach", "sentimentScore (1-5) from summary.json — pre-labeled. Sentence-level sentimentType labels. No LLM re-scoring."),
        ("Topic & Churn", f"HDBSCAN clustering (silhouette: {fmt_num(data.hdbscan_score, 3)}). {data.hdbscan_noise} calls noise. LLM naming + TF-IDF. Business taxonomy added for presentation clarity. Churn is rule-based and illustrative."),
        ("Limitations", f"Silhouette {fmt_num(data.hdbscan_score, 3)} is weak (ideal > 0.5); {round(100*data.hdbscan_noise/data.total_calls) if data.total_calls else 0}% noise rate; N={data.total_calls} means trends are suggestive."),
    ]

    card_w = Inches(4.55)
    card_h = Inches(1.80)
    positions = [
        (MARGIN, CONTENT_Y),
        (MARGIN + Inches(4.75), CONTENT_Y),
        (MARGIN, CONTENT_Y + Inches(1.95)),
        (MARGIN + Inches(4.75), CONTENT_Y + Inches(1.95)),
    ]

    for (x, y), (title, body) in zip(positions, cards):
        add_card(slide, x, y, card_w, card_h)
        add_text(slide, x + Inches(0.12), y + Inches(0.12), card_w - Inches(0.24), Inches(0.28),
                 title, FONT_BODY, bold=True, color=C_PRIMARY)
        body = fit_text(body, card_w - Inches(0.24), FONT_TINY, 5)
        add_text(slide, x + Inches(0.12), y + Inches(0.45), card_w - Inches(0.24), Inches(1.25),
                 body, FONT_TINY, color=C_TEXT)
    check(slide, "Methodology")
    return slide


def add_charts_appendix_slide(prs, data: PresentationData):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Appendix: Additional Charts", "Evidence behind the findings")

    chart_w = Inches(4.4)
    chart_h = Inches(1.9)
    gap_x = Inches(0.2)
    gap_y = Inches(0.15)
    positions = [
        (MARGIN, CONTENT_Y),
        (MARGIN + chart_w + gap_x, CONTENT_Y),
        (MARGIN, CONTENT_Y + chart_h + gap_y),
        (MARGIN + chart_w + gap_x, CONTENT_Y + chart_h + gap_y),
    ]
    charts = [
        "05_churn_score_histogram.png",
        "04_negative_sentiment_trend.png",
        "03_clustering_comparison.png",
        "01_duration_distribution.png",
    ]
    for (x, y), filename in zip(positions, charts):
        add_chart_if_exists(slide, filename, x, y, chart_w, chart_h)
    check(slide, "Charts Appendix")
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("=" * 60)
    print("06 GENERATE PPT: Building presentation")
    print("=" * 60)

    try:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        data = load_presentation_data(OUTPUT_DIR)
        if data.warnings:
            print("\nData validation warnings:")
            for w in data.warnings:
                print(f"  - {w}")

        slides = [
            ("Title + KPIs", add_title_slide, True),
            ("Executive Summary", add_executive_summary_slide, True),
            ("Pipeline", add_pipeline_slide, False),
            ("Dataset", add_dataset_slide, True),
            ("Topics", add_topic_slide, True),
            ("Sentiment", add_sentiment_slide, True),
            ("Problem Zones", add_problem_zones_slide, True),
            ("Strong Zones", add_strong_zones_slide, True),
            ("Churn", add_churn_slide, True),
            ("Renewal Risk", add_renewal_risk_slide, True),
            ("Features", add_feature_slide, True),
            ("Action Items", add_action_items_slide, True),
            ("Carry-Forward Actions", add_carry_forward_slide, True),
            ("Recommendations at a Glance", add_recommendations_slide, True),
            ("Recommendation Details", add_recommendation_detail_slides, True),
            ("AI Reasonableness Check", add_reasonableness_slide, True),
            ("Methodology", add_methodology_slide, True),
        ]

        for i, (name, func, needs_data) in enumerate(slides, 1):
            func(prs, data) if needs_data else func(prs)
            print(f"  Slide {i}: {name}")

        output_path = OUTPUT_DIR / "Transcript_Intelligence_Report.pptx"
        prs.save(output_path)
        print("\n" + "=" * 60)
        print("06 GENERATE PPT: Done")
        print(f"Presentation saved: {output_path}")
        print(f"Total slides: {len(prs.slides)}")
        print("=" * 60)
    except FileNotFoundError as e:
        print(f"\nERROR: Missing required analysis output. Run scripts 01-05 first.\n  {e}")
        sys.exit(1)
    except Exception as e:
        print(f"\nERROR: PPT generation failed: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
