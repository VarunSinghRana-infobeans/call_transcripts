"""
ppt_to_md.py

Convert the generated PowerPoint deck to Markdown for review.
Extracts slide titles, body text, tables, and notes a rough layout map.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def shape_type_name(shape):
    """Return a human-readable shape category."""
    if shape.has_text_frame:
        if shape.text_frame.text.strip():
            return "text"
    if shape.shape_type == 13:  # PICTURE
        return "image"
    if shape.has_table:
        return "table"
    return "shape"


def describe_shape(shape, slide_idx):
    """Return a Markdown snippet describing one shape."""
    parts = []
    left = shape.left.inches
    top = shape.top.inches
    width = shape.width.inches
    height = shape.height.inches

    stype = shape_type_name(shape)

    if stype == "text":
        tf = shape.text_frame
        text = tf.text.strip()
        if not text:
            return None
        # Estimate font size from first run
        font_size = None
        for p in tf.paragraphs:
            for run in p.runs:
                if run.font.size:
                    font_size = run.font.size.pt
                    break
            if font_size:
                break
        bold = any(run.font.bold for p in tf.paragraphs for run in p.runs if run.font.bold is not None)
        size_hint = f" (~{font_size:.0f}pt)" if font_size else ""
        pos = f"[{left:.2f}\", {top:.2f}\"] {width:.2f}\"×{height:.2f}\""
        lines = [f"- **Text box** {pos}{' bold' if bold else ''}{size_hint}"]
        for p in tf.paragraphs:
            line = p.text.strip()
            if line:
                lines.append(f"  - {line}")
        return "\n".join(lines)

    if stype == "image":
        return f"- **Image** at [{left:.2f}\", {top:.2f}\"] {width:.2f}\"×{height:.2f}\""

    if stype == "table":
        tbl = shape.table
        rows = len(tbl.rows)
        cols = len(tbl.columns)
        return f"- **Table** {rows}×{cols} at [{left:.2f}\", {top:.2f}\"] {width:.2f}\"×{height:.2f}\""

    return f"- **{stype}** at [{left:.2f}\", {top:.2f}\"] {width:.2f}\"×{height:.2f}\""


def slide_to_md(slide, idx, slide_width_in, slide_height_in):
    """Return Markdown for one slide."""
    # Try to find a title: top-most text box near the top
    title = None
    title_candidates = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                title_candidates.append((shape.top.inches, text, shape))
    title_candidates.sort(key=lambda x: x[0])
    if title_candidates and title_candidates[0][0] < 1.2:
        title = title_candidates[0][1]
        title_shape = title_candidates[0][2]
        title_fs = None
        for p in title_shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    title_fs = run.font.size.pt
                    break
            if title_fs:
                break
        title_size = f" (~{title_fs:.0f}pt)" if title_fs else ""
    else:
        title = f"Slide {idx}"
        title_size = ""

    md = [f"## Slide {idx}: {title}{title_size}\n"]
    md.append(f"Canvas: {slide_width_in:.2f}\" × {slide_height_in:.2f}\"\n")

    # Sort shapes top-to-bottom, then left-to-right
    shapes_info = []
    for shape in slide.shapes:
        desc = describe_shape(shape, idx)
        if desc:
            shapes_info.append((shape.top.inches, shape.left.inches, desc))
    shapes_info.sort(key=lambda x: (x[0], x[1]))

    for _, _, desc in shapes_info:
        md.append(desc)

    md.append("")
    return "\n".join(md)


def ppt_to_md(ppt_path: Path, out_path: Path | None = None) -> Path:
    ppt_path = Path(ppt_path)
    prs = Presentation(str(ppt_path))
    width = prs.slide_width.inches

    md_lines = [
        f"# {ppt_path.name}\n",
        f"Slides: {len(prs.slides)} | Canvas: {width:.2f}\" × {prs.slide_height.inches:.2f}\"\n",
    ]

    height = prs.slide_height.inches
    for idx, slide in enumerate(prs.slides, start=1):
        md_lines.append(slide_to_md(slide, idx, width, height))

    md_text = "\n".join(md_lines)

    if out_path is None:
        out_path = ppt_path.with_suffix(".md")
    out_path.write_text(md_text, encoding="utf-8")
    print(f"Markdown saved: {out_path}")
    return out_path


if __name__ == "__main__":
    if len(sys.argv) > 1:
        path = Path(sys.argv[1])
    else:
        path = Path(__file__).parent / "output" / "latest" / "Transcript_Intelligence_Report.pptx"
    ppt_to_md(path)
